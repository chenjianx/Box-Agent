"""Regression tests for PPTX creation guardrails in execute_code."""

from __future__ import annotations

import pytest

from box_agent.tools.file_tools import MAX_FILE_TOOL_CONTENT_CHARS
from box_agent.tools.jupyter_tool import MAX_EXECUTE_CODE_CHARS, JupyterSandboxTool


def test_execute_code_schema_exposes_code_size_limit():
    code_schema = JupyterSandboxTool().parameters["properties"]["code"]

    assert code_schema["maxLength"] == MAX_EXECUTE_CODE_CHARS
    assert f"{MAX_EXECUTE_CODE_CHARS:,} characters" in code_schema["description"]
    assert "do not inline the body in execute_code" in code_schema["description"]
    assert "ordered write_file chunks" in code_schema["description"]
    assert "chunk_index/final" in code_schema["description"]
    assert "JSON manifests" in code_schema["description"]


def test_execute_code_limit_is_not_looser_than_file_chunk_limit():
    assert MAX_EXECUTE_CODE_CHARS <= MAX_FILE_TOOL_CONTENT_CHARS


@pytest.mark.asyncio
async def test_execute_code_rejects_oversized_code_before_kernel_start(monkeypatch):
    def fail_if_sandbox_requested(self):
        raise AssertionError("oversized code should be rejected before sandbox startup")

    monkeypatch.setattr(JupyterSandboxTool, "_get_sandbox_env", fail_if_sandbox_requested)
    tool = JupyterSandboxTool()
    code = "x = 1\n" + ("print(x)\n" * MAX_EXECUTE_CODE_CHARS)

    result = await tool.execute(code=code)

    assert result.success is False
    assert result.error is not None
    assert result.error.startswith("EXECUTE_CODE_TOO_LARGE")
    assert "Split the work into multiple execute_code calls" in result.error
    assert "ordered write_file chunks" in result.error
    assert "chunk_index/final" in result.error


def test_execute_code_blocks_bare_python_pptx_new_deck_constructor():
    code = """\
from pptx import Presentation

prs = Presentation()
prs.save("deck.pptx")
"""

    assert JupyterSandboxTool._looks_like_python_pptx_new_deck(code)


def test_execute_code_allows_existing_pptx_inspection():
    code = """\
from pptx import Presentation

prs = Presentation("existing.pptx")
print(len(prs.slides))
"""

    assert not JupyterSandboxTool._looks_like_python_pptx_new_deck(code)


def test_execute_code_allows_non_pptx_data_work():
    code = """\
import pandas as pd

df = pd.read_csv("data.csv")
print(df.describe())
"""

    assert not JupyterSandboxTool._looks_like_python_pptx_new_deck(code)


def test_execute_code_detects_controlled_deck_json_rewrite_but_allows_read() -> None:
    rewrite = """\
import json
with open("/tmp/output/deck.json", "w") as handle:
    json.dump(deck, handle)
"""
    read_only = """\
import json
with open("/tmp/output/deck.json") as handle:
    deck = json.load(handle)
print(len(deck["slides"]))
"""

    assert JupyterSandboxTool._looks_like_controlled_deck_rewrite(rewrite)
    assert not JupyterSandboxTool._looks_like_controlled_deck_rewrite(read_only)


@pytest.mark.parametrize(
    "rewrite",
    [
        'from pathlib import Path\nPath("deck.json").open("w").write("{}")',
        'from pathlib import Path\ndeck_path = Path("deck.json")\ndeck_path.write_text("{}")',
        'from pathlib import Path\ndeck_path = Path("output") / "deck.json"\ndeck_path.write_text("{}")',
        'deck_path = "/tmp/output/deck.json"\nopen(deck_path, "w").write("{}")',
    ],
)
def test_execute_code_detects_indirect_controlled_deck_rewrites(rewrite: str) -> None:
    assert JupyterSandboxTool._looks_like_controlled_deck_rewrite(rewrite)


@pytest.mark.asyncio
async def test_execute_code_blocks_controlled_deck_rewrite_before_kernel_start(
    monkeypatch,
) -> None:
    def fail_if_sandbox_requested(self):
        raise AssertionError("deck rewrite should be rejected before sandbox startup")

    monkeypatch.setattr(JupyterSandboxTool, "_get_sandbox_env", fail_if_sandbox_requested)
    tool = JupyterSandboxTool()

    result = await tool.execute(
        code='from pathlib import Path\nPath("deck.json").write_text("{}")'
    )

    assert result.success is False
    assert result.error is not None
    assert result.error.startswith("CONTROLLED_DECK_REWRITE_BLOCKED")
    assert "apply_deck_patch.js" in result.error
